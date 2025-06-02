"""
DocToPPT - Aplicação Flask Principal
Gerador de Apresentações PowerPoint com IA
"""

import os
import logging
from flask import Flask, render_template, request, jsonify, send_file, flash, redirect, url_for
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
from config import config
import threading
import time
import shutil
from utils import check_and_install_missing_packages, extract_pdf_images

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def create_app(config_name=None):
    """Application factory pattern"""
    if config_name is None:
        config_name = os.getenv('FLASK_ENV', 'default')
    
    app = Flask(__name__)
    app.config.from_object(config[config_name])
    config[config_name].init_app(app)
    
    # Função para processamento de documentos em background
    def process_document_async(file_path, output_filename, template_path=None):
        """
        Função que processa um documento e cria um PowerPoint
        Extrai o texto do PDF e cria slides básicos
        Parâmetros:
          file_path: Caminho do arquivo de entrada (PDF, DOCX, etc)
          output_filename: Nome do arquivo de saída (sem extensão)
          template_path: Caminho para um arquivo PPTX de template (opcional)
        """
        logger.info(f"Iniciando processamento assíncrono para {file_path}")
        
        # Importar bibliotecas no escopo da função para não afetar o tempo de inicialização
        try:
            import PyPDF2
            import pdfplumber
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            import re
        except ImportError as e:
            logger.error(f"Erro ao importar bibliotecas: {e}")
            return False
            
        # Criar um arquivo de saída
        output_dir = os.path.join(app.root_path, 'static', 'outputs')
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"{output_filename}.pptx")
        
        try:
            # Verificar a extensão do arquivo para saber como processá-lo
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # Extrair texto do documento baseado no tipo
            text_content = []
            page_metrics = []  # Armazenar métricas de cada página para análise
            
            if file_ext == '.pdf':
                # Extrair texto do PDF usando pdfplumber para melhor extração de estrutura
                image_paths = []  # Lista para armazenar caminhos das imagens extraídas
                
            # Tentar extrair imagens do PDF
                try:
                    # Criar diretório temporário para imagens extraídas
                    image_dir = os.path.join(output_dir, f"images_{output_filename}")
                    os.makedirs(image_dir, exist_ok=True)
                    
                    # Chamar função para extrair imagens - usando a função definida neste escopo
                    image_paths = extract_pdf_images(file_path, image_dir)
                    if image_paths:
                        logger.info(f"Extraídas {len(image_paths)} imagens do PDF")
                except Exception as img_ex:
                    logger.error(f"Erro ao extrair imagens: {img_ex}")
                
                try:
                    with pdfplumber.open(file_path) as pdf:
                        num_pages = len(pdf.pages)
                        logger.info(f"PDF tem {num_pages} páginas")
                        
                        # Extrair texto de cada página com preservação de estrutura
                        for page_num in range(min(num_pages, 20)):  # Limitar a 20 páginas
                            page = pdf.pages[page_num]
                            
                            # Extrair métricas de texto para detectar cabeçalhos
                            page_text = ""
                            structured_lines = []
                            
                            # Extrair texto com informações de fonte
                            text_objects = page.extract_words(x_tolerance=3, y_tolerance=3, keep_blank_chars=False, 
                                                           use_text_flow=True, extra_attrs=['size', 'fontname'])
                            
                            # Agrupar palavras em linhas preservando informações de tamanho da fonte
                            current_y = None
                            current_line = []
                            current_size = None
                            current_font = None
                            
                            for word in text_objects:
                                if current_y is None or abs(word['top'] - current_y) > 5:
                                    # Nova linha detectada
                                    if current_line:
                                        line_text = " ".join([w['text'] for w in current_line])
                                        structured_lines.append({
                                            'text': line_text,
                                            'size': current_size,
                                            'font': current_font,
                                            'y': current_y,
                                            'page': page_num
                                        })
                                        page_text += line_text + "\n"
                                    
                                    current_line = [word]
                                    current_y = word['top']
                                    current_size = word['size']
                                    current_font = word['fontname']
                                else:
                                    # Continuar mesma linha
                                    current_line.append(word)
                                    # Atualizar tamanho de fonte se for maior (provável destaque)
                                    if word['size'] > current_size:
                                        current_size = word['size']
                            
                            # Adicionar última linha
                            if current_line:
                                line_text = " ".join([w['text'] for w in current_line])
                                structured_lines.append({
                                    'text': line_text,
                                    'size': current_size,
                                    'font': current_font,
                                    'y': current_y,
                                    'page': page_num
                                })
                                page_text += line_text + "\n"
                            
                            text_content.append(page_text)
                            page_metrics.append(structured_lines)
                            
                            # Log para debug
                            if page_num == 0:
                                preview = page_text[:200] + "..." if len(page_text) > 200 else page_text
                                logger.info(f"Amostra de texto extraído com pdfplumber: {preview}")
                                logger.info(f"Estrutura de fonte identificada: {len(structured_lines)} linhas")
                    
                    # Adicionar informação sobre imagens encontradas aos metadados
                    if image_paths:
                        # Mapear imagens para páginas com base nos nomes de arquivo
                        images_by_page = {}
                        for img_path in image_paths:
                            img_name = os.path.basename(img_path)
                            # Extrair número da página do nome da imagem (formato: image_p{page_num}_{img_index}.{ext})
                            match = re.search(r'image_p(\d+)_', img_name)
                            if match:
                                page_num = int(match.group(1)) - 1  # Ajustar para base 0
                                if page_num not in images_by_page:
                                    images_by_page[page_num] = []
                                images_by_page[page_num].append(img_path)
                        
                        # Adicionar metadados de imagens à estrutura
                        for page_num, page_images in images_by_page.items():
                            if page_num < len(text_content):
                                text_content[page_num] += f"\n[{len(page_images)} imagens encontradas nesta página]\n"
                                
                                # Associar imagens com texto próximo
                                if page_num < len(page_metrics):
                                    page_metrics[page_num].append({
                                        'text': f"[IMAGENS: {len(page_images)}]",
                                        'size': 12,
                                        'font': 'Image',
                                        'y': 999999,  # Valor alto para ser processado por último
                                        'page': page_num,
                                        'images': page_images
                                    })
                except Exception as pdfex:
                    logger.error(f"Erro ao processar PDF com pdfplumber: {pdfex}")
                    logger.info("Tentando método alternativo com PyPDF2...")
                    
                    # Método alternativo com PyPDF2 caso pdfplumber falhe
                    with open(file_path, 'rb') as pdf_file:
                        reader = PyPDF2.PdfReader(pdf_file)
                        num_pages = len(reader.pages)
                        
                        # Extrair texto de cada página
                        for page_num in range(min(num_pages, 20)):  # Limitar a 20 páginas
                            page = reader.pages[page_num]
                            extracted_text = page.extract_text() or ""
                            text_content.append(extracted_text)
                            
                            # Para compatibilidade com o pipeline
                            lines = extracted_text.split('\n')
                            structured_lines = [{'text': line, 'size': 12, 'font': 'Unknown', 'page': page_num} for line in lines]
                            
                            # Adicionar informações sobre imagens, se existirem
                            if page_num in images_by_page:
                                page_images = images_by_page[page_num]
                                extracted_text += f"\n[{len(page_images)} imagens encontradas]\n"
                                structured_lines.append({
                                    'text': f"[IMAGENS: {len(page_images)}]",
                                    'size': 12,
                                    'font': 'Image',
                                    'page': page_num,
                                    'images': page_images
                                })
                            
                            page_metrics.append(structured_lines)
                            
                            # Log para debug
                            if page_num == 0:
                                preview = extracted_text[:200] + "..." if len(extracted_text) > 200 else extracted_text
                                logger.info(f"Amostra de texto extraído com PyPDF2: {preview}")
            
            elif file_ext in ['.docx', '.txt', '.md']:
                # Para outros formatos, ler como texto simples (implementação básica)
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as text_file:
                    text_content.append(text_file.read())
            else:
                logger.error(f"Formato de arquivo não suportado: {file_ext}")
                return False
            
            # Processar o texto extraído para criar slides
            # Aqui seria o lugar para chamar a API de IA para organizar o conteúdo
            # Por enquanto, vamos fazer uma implementação simples dividindo o texto em partes
            
            # Criar uma apresentação (baseada em template se fornecido)
            if template_path and os.path.exists(template_path):
                logger.info(f"Usando template: {template_path}")
                try:
                    prs = Presentation(template_path)
                    # Usar template existente
                    logger.info(f"Template carregado com {len(prs.slides)} slides existentes")
                except Exception as te:
                    logger.error(f"Erro ao carregar template: {str(te)}. Usando apresentação padrão.")
                    prs = Presentation()
            else:
                prs = Presentation()
            
            # Slide de título (adicionar ou usar o primeiro slide existente)
            if len(prs.slides) > 0 and template_path:
                # Usar o primeiro slide do template como título
                slide = prs.slides[0]
                try:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            if shape == slide.shapes.title:
                                shape.text = f"{output_filename} - Apresentação"
                            else:
                                shape.text = "Gerado automaticamente por DocToPPT"
                            break
                except Exception as e:
                    logger.warning(f"Não foi possível editar o slide de título do template: {str(e)}")
            else:
                # Criar slide de título
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = f"{output_filename} - Apresentação"
                subtitle.text = "Gerado automaticamente por DocToPPT"
            
            # Adicionar slides com o conteúdo extraído
            bullet_slide_layout = prs.slide_layouts[1]
            
            # Processar o texto extraído e separar em partes significativas usando
            # informações de tamanho de fonte e estrutura detectada
            sections = []
            current_section = {"title": "Introdução", "content": []}
            
            # Analisar as métricas para identificar padrões de cabeçalhos
            if page_metrics and len(page_metrics) > 0:
                # Identificar tamanhos de fonte comuns para uso na detecção de cabeçalhos
                all_sizes = []
                for page in page_metrics:
                    for line in page:
                        if 'size' in line and line['size']:
                            all_sizes.append(line['size'])
                
                if all_sizes:
                    # Calcular estatísticas sobre tamanhos de fonte
                    avg_size = sum(all_sizes) / len(all_sizes)
                    larger_sizes = sorted([s for s in all_sizes if s > avg_size])
                    
                    # Definir limites para cabeçalhos baseados na análise estatística
                    heading_size_threshold = avg_size * 1.1 if larger_sizes else avg_size
                    logger.info(f"Tamanho médio de fonte: {avg_size:.2f}, limiar para cabeçalhos: {heading_size_threshold:.2f}")
                    
                    # Padrões regulares que indicam cabeçalhos
                    heading_patterns = [
                        r'^[0-9]+\.\s+',           # Números de seção: "1. Título"
                        r'^[A-Z\s]{5,}$',          # Texto todo em maiúsculas
                        r'^(CAPÍTULO|SEÇÃO|PARTE)', # Palavras específicas
                        r'^[IVX]+\.\s+'            # Numerais romanos: "IV. Título"
                    ]
                    
                    # Processar as páginas com métricas
                    page_num = 0
                    for page in page_metrics:
                        # Adicionar marcador de página para melhor segmentação
                        if page_num > 0 and current_section["content"]:
                            current_section["content"].append("\n[Nova Página]\n")
                        
                        for line in page:
                            text = line['text'].strip()
                            if not text:
                                continue
                                
                            is_heading = False
                            
                            # Verificar pelo tamanho da fonte
                            if 'size' in line and line['size'] and line['size'] > heading_size_threshold:
                                is_heading = True
                                logger.debug(f"Cabeçalho detectado por tamanho de fonte: {text} ({line['size']})")
                            
                            # Verificar pelos padrões de texto
                            if not is_heading and len(text) < 80:  # Cabeçalhos geralmente são curtos
                                for pattern in heading_patterns:
                                    if re.match(pattern, text):
                                        is_heading = True
                                        logger.debug(f"Cabeçalho detectado por padrão: {text}")
                                        break
                            
                            # Verificar outros indicadores (numeração, etc)
                            if not is_heading and len(text) < 60:
                                if (text.isupper() or text.endswith(':') or
                                    any(text.startswith(prefix) for prefix in ['CAPÍTULO', 'SEÇÃO', 'INTRODUÇÃO', 'CONCLUSÃO'])):
                                    is_heading = True
                            
                            # Tratar como cabeçalho ou conteúdo
                            if is_heading:
                                # Começar nova seção se a atual já tem conteúdo
                                if current_section["content"]:
                                    sections.append(current_section)
                                    current_section = {"title": text, "content": []}
                                else:
                                    # Atualizar o título da seção atual se ainda não tem conteúdo
                                    current_section["title"] = text
                            else:
                                # Adicionar ao conteúdo da seção atual
                                current_section["content"].append(text)
                        
                        page_num += 1
                        
                # Log da estrutura detectada
                logger.info(f"Estrutura detectada: {len(sections) + 1} seções")
            else:
                # Fallback para o método anterior de detecção simples de seções
                logger.warning("Usando método alternativo para detecção de seções")
                for page_text in text_content:
                    lines = page_text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                            
                        # Detectar possíveis títulos (linhas curtas, todas em maiúsculas, etc.)
                        if (len(line) < 60 and (line.isupper() or line.endswith(':') or 
                                            line.startswith('CAPÍTULO') or line.startswith('SEÇÃO'))):
                            # Começar nova seção se a atual já tem conteúdo
                            if current_section["content"]:
                                sections.append(current_section)
                                current_section = {"title": line, "content": []}
                        else:
                            # Adicionar linha ao conteúdo da seção atual
                            current_section["content"].append(line)
            
            # Adicionar a última seção
            if current_section["content"]:
                sections.append(current_section)
            
            # Se não conseguimos identificar seções, criar uma divisão artificial
            if not sections:
                sections = [{"title": "Conteúdo", "content": " ".join(text_content)}]
                
            # Dividir em chunks se o conteúdo for muito grande
            processed_sections = []
            for section in sections:
                content = section["content"]
                if isinstance(content, list):
                    content = " ".join(content)
                    
                # Dividir em chunks para não sobrecarregar os slides
                if len(content) > 800:
                    chunks = [content[i:i+800] for i in range(0, len(content), 800)]
                    for i, chunk in enumerate(chunks):
                        title = section["title"]
                        if i > 0:
                            title += f" (cont. {i+1})"
                        processed_sections.append({"title": title, "content": chunk})
                else:
                    processed_sections.append(section)
            
            # Determinar qual layout usar para os slides de conteúdo
            # Preferir layouts do template, se disponível
            content_layouts = []
            if template_path and len(prs.slide_layouts) > 1:
                # Usar layouts do template para os slides de conteúdo
                # Coletar layouts que parecem adequados para conteúdo
                for layout in prs.slide_layouts:
                    # Verificar se o layout tem placeholder para título e conteúdo
                    has_title = False
                    has_content = False
                    for placeholder in layout.placeholders:
                        if placeholder.placeholder_format.type == 1:  # TITLE
                            has_title = True
                        if placeholder.placeholder_format.type in [2, 7]:  # BODY, CONTENT
                            has_content = True
                    
                    if has_title and has_content:
                        content_layouts.append(layout)
            
            # Se não encontrou layouts adequados no template, usar o padrão
            if not content_layouts:
                content_layouts = [prs.slide_layouts[1]]  # Layout de título e conteúdo padrão
            
            # Criar slides baseados nas seções processadas
            for i, section in enumerate(processed_sections[:20]):  # Limitar a 20 slides
                # Verificar se a seção contém imagens
                has_images = False
                section_images = []
                
                # Procurar por imagens no conteúdo
                if isinstance(section["content"], list):
                    for item in section["content"]:
                        if isinstance(item, str) and "[IMAGENS:" in item:
                            has_images = True
                        if isinstance(item, dict) and "images" in item:
                            has_images = True
                            section_images.extend(item["images"])
                            
                # Escolher layout apropriado
                layout_index = i % len(content_layouts)
                
                # Se tiver imagens, tentar encontrar um layout com placeholder para imagens
                # Por enquanto, vamos usar layouts padrão pois não sabemos quais têm placeholder para imagem
                layout = content_layouts[layout_index]
                slide = prs.slides.add_slide(layout)
                
                # Encontrar placeholder para título
                title = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type == 1:  # TITLE
                        title = shape
                        break
                
                # Encontrar placeholder para conteúdo
                content = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in [2, 7]:  # BODY, CONTENT
                        content = shape
                        break
                
                # Se não encontrou os placeholders esperados, pular este slide
                if not title or not content:
                    logger.warning(f"Layout {layout_index} não tem placeholders esperados, pulando.")
                    continue
                
                # Adicionar as imagens encontradas nesta seção
                if section_images:
                    # Adicionar até 2 imagens por slide
                    images_to_add = section_images[:2]
                    
                    try:
                        # Calcular posições para as imagens
                        if len(images_to_add) == 1:
                            # Uma imagem - posicionar à direita
                            img_path = images_to_add[0]
                            # Verificar se o arquivo existe
                            if os.path.exists(img_path):
                                # Adicionar imagem ao slide
                                try:
                                    left = Inches(5)  # Posição à direita
                                    top = Inches(2)
                                    width = Inches(4)  # Largura menor para não cobrir texto
                                    slide.shapes.add_picture(img_path, left, top, width=width)
                                    logger.info(f"Imagem adicionada ao slide {i+1}: {img_path}")
                                except Exception as img_err:
                                    logger.error(f"Erro ao adicionar imagem: {img_err}")
                        else:
                            # Múltiplas imagens - distribuir
                            for idx, img_path in enumerate(images_to_add):
                                if os.path.exists(img_path):
                                    try:
                                        # Posição ajustada para cada imagem
                                        left = Inches(1 + (idx * 3))  # Espaçamento horizontal
                                        top = Inches(4)  # Abaixo do conteúdo
                                        width = Inches(3)  # Largura menor
                                        slide.shapes.add_picture(img_path, left, top, width=width)
                                        logger.info(f"Imagem {idx+1} adicionada ao slide {i+1}: {img_path}")
                                    except Exception as img_err:
                                        logger.error(f"Erro ao adicionar imagem {idx+1}: {img_err}")
                    except Exception as ex:
                        logger.error(f"Erro ao processar imagens para slide {i+1}: {ex}")
                
                title.text = section["title"]
                
                # Limitar o tamanho do texto para não sobrecarregar o slide
                text_content = section["content"]
                
                # Converter para texto simples se for lista
                if isinstance(text_content, list):
                    # Processa o conteúdo para detectar listas e formatação
                    formatted_content = []
                    bullet_mode = False
                    
                    for line in text_content:
                        # Detectar marcadores de lista
                        if re.match(r'^[\s]*[•\-\*\+◦○●■]\s+', line) or re.match(r'^[\s]*[0-9]+[\.\)]\s+', line):
                            if not bullet_mode:
                                # Iniciar modo de lista
                                bullet_mode = True
                                if formatted_content and formatted_content[-1].strip():
                                    formatted_content.append('')  # Linha em branco antes da lista
                            # Remover o marcador e adicionar como item de lista com formato apropriado
                            clean_line = re.sub(r'^[\s]*[•\-\*\+◦○●■]\s+', '', line)
                            clean_line = re.sub(r'^[\s]*[0-9]+[\.\)]\s+', '', clean_line)
                            formatted_content.append('• ' + clean_line.strip())
                        else:
                            # Terminar modo de lista se estiver ativo
                            if bullet_mode and line.strip():
                                bullet_mode = False
                                formatted_content.append('')  # Linha em branco depois da lista
                            formatted_content.append(line)
                    
                    # Juntar com quebras de linha
                    text_content = "\n".join(formatted_content)
                    
                # Truncar se for muito longo
                if len(text_content) > 1500:
                    text_content = text_content[:1500] + "..."
                
                # Adicionar ao texto frame com formatação
                tf = content.text_frame
                tf.clear()
                
                # Verificar se o conteúdo contém marcadores de lista
                if '• ' in text_content:
                    # Dividir por linhas e adicionar parágrafos com nível apropriado
                    lines = text_content.split('\n')
                    first = True
                    for line in lines:
                        if first:
                            p = tf.paragraphs[0]
                            first = False
                        else:
                            p = tf.add_paragraph()
                            
                        # Detectar se é um item de lista
                        if line.startswith('• '):
                            p.text = line[2:]  # Remover o marcador
                            p.level = 1  # Nível de recuo para item de lista
                        else:
                            p.text = line
                            p.level = 0  # Nível normal para texto
                else:
                    # Texto simples sem formatação especial
                    tf.text = text_content
            
            # Adicionar slide de conclusão
            conclusion_layout = content_layouts[0]  # Usar primeiro layout de conteúdo
            slide = prs.slides.add_slide(conclusion_layout)
            
            # Encontrar placeholders para título e conteúdo
            title = None
            content = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # TITLE
                    title = shape
                elif shape.placeholder_format.type in [2, 7]:  # BODY, CONTENT
                    content = shape
            
            if title and content:
                title.text = "Conclusão"
                content.text = "Obrigado!\n\nEste documento foi gerado automaticamente pelo DocToPPT."
            
            # Salvar a apresentação
            prs.save(output_path)
            
            logger.info(f"Processamento concluído. Arquivo PowerPoint criado: {output_path}")
            return True
        except Exception as e:
            logger.error(f"Erro durante o processamento de documento: {str(e)}")
            
            # Criar arquivo de PowerPoint de erro para notificar o usuário
            try:
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                title.text = "Erro no Processamento"
                subtitle.text = f"Ocorreu um erro ao processar o documento: {str(e)}"
                
                prs.save(output_path)
                logger.info(f"Arquivo de erro criado: {output_path}")
                return True
            except Exception as e2:
                logger.error(f"Falha ao criar arquivo de erro: {str(e2)}")
                return False
    
    # Create upload directory if it doesn't exist
    upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
    upload_dir = os.path.join(app.root_path, upload_folder)
    os.makedirs(upload_dir, exist_ok=True)
    logger.info(f"Upload directory configured: {upload_dir}")
    
    # Create output directory for generated presentations
    output_dir = os.path.join(app.root_path, 'static', 'outputs')
    os.makedirs(output_dir, exist_ok=True)
    logger.info(f"Output directory configured: {output_dir}")
    
    # Error handlers
    @app.errorhandler(RequestEntityTooLarge)
    def handle_file_too_large(e):
        flash('Arquivo muito grande. Tamanho máximo: 16MB', 'error')
        return redirect(url_for('index')), 413
    
    @app.errorhandler(404)
    def not_found(e):
        return render_template('error.html', error="Página não encontrada"), 404
    
    @app.errorhandler(500)
    def server_error(e):
        logger.error(f"Server error: {e}")
        return render_template('error.html', error="Erro interno do servidor"), 500
    
    # Utility functions
    def allowed_file(filename):
        """Check if file extension is allowed"""
        if '.' not in filename:
            return False
        
        ext = filename.rsplit('.', 1)[1].lower()
        allowed = ext in app.config['ALLOWED_EXTENSIONS']
        
        logger.info(f"File extension check: {filename} -> {ext}, allowed: {allowed}")
        return allowed
    
    # Routes
    @app.route('/')
    def index():
        """Home page"""
        return render_template('index.html')
    
    @app.route('/upload', methods=['GET', 'POST'])
    def upload():
        """Upload page and file handling"""
        if request.method == 'GET':
            return render_template('upload.html')
        
        # Log received form data
        logger.info(f"Upload request received. Content-Type: {request.content_type}")
        logger.info(f"Form data: {list(request.form.keys()) if request.form else 'No form data'}")
        logger.info(f"Files: {list(request.files.keys()) if request.files else 'No files'}")
        
        if 'document' not in request.files:
            logger.error("No document part in the request")
            flash('Nenhum arquivo selecionado', 'error')
            return redirect(request.url)
        
        file = request.files['document']
        template_file = request.files.get('template')
        
        if file.filename == '':
            logger.error("Empty filename submitted")
            flash('Nenhum arquivo selecionado', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename):
            try:
                filename = secure_filename(file.filename)
                logger.info(f"Arquivo seguro: {filename}")
                
                # Garantir que o diretório de upload existe
                upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
                upload_dir = os.path.join(app.root_path, upload_folder)
                os.makedirs(upload_dir, exist_ok=True)
                logger.info(f"Diretório de upload: {upload_dir}")
                
                # Salvar o arquivo com caminho completo
                file_path = os.path.join(upload_dir, filename)
                logger.info(f"Tentando salvar arquivo em: {file_path}")
                try:
                    file.save(file_path)
                    logger.info(f"File uploaded successfully: {filename} to {file_path}")
                    
                    # Verificar se o arquivo foi realmente salvo
                    if os.path.exists(file_path):
                        logger.info(f"Confirmado: arquivo existe no caminho {file_path}")
                        flash(f'Arquivo {filename} enviado com sucesso!', 'success')
                        
                        # Verificar se o usuário enviou um template
                        template_path = None
                        if template_file and template_file.filename:
                            template_filename = secure_filename(template_file.filename)
                            template_path = os.path.join(upload_dir, f"template_{template_filename}")
                            try:
                                template_file.save(template_path)
                                logger.info(f"Template salvo: {template_path}")
                            except Exception as te:
                                logger.error(f"Erro ao salvar template: {str(te)}")
                                template_path = None
                            
                        # Iniciar processamento em segundo plano
                        output_filename = os.path.splitext(filename)[0]
                        threading.Thread(
                            target=process_document_async, 
                            args=(file_path, output_filename, template_path)
                        ).start()
                        
                        return redirect(url_for('processing', filename=filename))
                    else:
                        logger.error(f"Arquivo não encontrado após save: {file_path}")
                        flash('Erro: Arquivo não foi salvo corretamente', 'error')
                        return redirect(request.url)
                except IOError as io_err:
                    logger.error(f"IOError ao salvar arquivo: {io_err}")
                    flash(f'Erro de E/S ao salvar arquivo: {str(io_err)}', 'error')
                    return redirect(request.url)
            except Exception as e:
                logger.error(f"Error saving file: {e}")
                flash(f'Erro ao salvar arquivo: {str(e)}', 'error')
                return redirect(request.url)
        else:
            flash('Tipo de arquivo não permitido', 'error')
            return redirect(request.url)
    
    @app.route('/processing/<filename>')
    def processing(filename):
        """Processing status page"""
        # Verificar se o arquivo existe
        upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
        upload_dir = os.path.join(app.root_path, upload_folder)
        file_path = os.path.join(upload_dir, secure_filename(filename))
        
        if not os.path.exists(file_path):
            flash('Arquivo não encontrado. Por favor faça upload novamente.', 'error')
            return redirect(url_for('upload'))
            
        # Se o arquivo existe, mostrar a página de processamento
        return render_template('processing.html', filename=filename)
    
    @app.route('/config')
    def config_page():
        """Configuration page"""
        current_config = {
            'deepseek_configured': bool(app.config.get('DEEPSEEK_API_KEY')),
            'max_slides': app.config['MAX_SLIDES'],
            'language': app.config['DEFAULT_LANGUAGE'],
            'ai_temperature': app.config['AI_TEMPERATURE']
        }
        return render_template('config.html', config=current_config)
    
    @app.route('/result/<filename>')
    def result(filename):
        """Result page"""
        secure_name = secure_filename(filename)
        
        # Verificar se o arquivo de entrada existe
        upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
        upload_dir = os.path.join(app.root_path, upload_folder)
        file_path = os.path.join(upload_dir, secure_name)
        
        if not os.path.exists(file_path):
            flash('Arquivo de entrada não encontrado. Por favor faça upload novamente.', 'error')
            return redirect(url_for('upload'))
        
        # Verificar se o arquivo de saída existe
        output_filename = os.path.splitext(secure_name)[0]
        output_path = os.path.join(app.root_path, 'static', 'outputs', f"{output_filename}.pptx")
        
        if not os.path.exists(output_path):
            logger.warning(f"Arquivo de saída não encontrado: {output_path}")
            flash('A apresentação ainda está sendo processada. Por favor aguarde.', 'warning')
            return redirect(url_for('processing', filename=filename))
        
        # Se o arquivo existe, mostrar a página de resultado
        output_url = url_for('static', filename=f'outputs/{output_filename}.pptx')
        logger.info(f"Mostrando resultado para {filename}, download em: {output_url}")
        
        return render_template('result.html', 
                               filename=filename,
                               output_url=output_url,
                               file_size=os.path.getsize(output_path))
    
    @app.route('/download/<filename>')
    def download(filename):
        """Download do arquivo PowerPoint gerado"""
        output_filename = os.path.splitext(secure_filename(filename))[0]
        output_path = os.path.join(app.root_path, 'static', 'outputs', f"{output_filename}.pptx")
        
        if not os.path.exists(output_path):
            flash('Arquivo não encontrado', 'error')
            return redirect(url_for('index'))
            
        # Enviar o arquivo para download
        return send_file(
            output_path, 
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=f"{output_filename}.pptx"
        )
    
    @app.route('/health')
    def health():
        """Health check endpoint for Docker"""
        return jsonify({
            'status': 'healthy',
            'version': '0.1.0',
            'deepseek_configured': bool(app.config.get('DEEPSEEK_API_KEY'))
        })
    
    @app.route('/api/status/<filename>')
    def api_status(filename):
        """API endpoint for processing status"""
        # Para demonstração, vamos implementar uma simulação de processamento real
        # Em produção, isso seria baseado em um sistema de filas ou base de dados
        import time
        
        # Verificamos se o arquivo de entrada existe
        upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
        upload_dir = os.path.join(app.root_path, upload_folder)
        
        # Caminho do arquivo de entrada
        secure_name = secure_filename(filename)
        file_path = os.path.join(upload_dir, secure_name)
        
        if not os.path.exists(file_path):
            return jsonify({
                'status': 'error',
                'progress': 0,
                'message': 'Arquivo de entrada não encontrado'
            })
        
        # Verifica se o processamento foi concluído (arquivo de saída existe)
        output_filename = os.path.splitext(secure_name)[0]
        output_path = os.path.join(app.root_path, 'static', 'outputs', f"{output_filename}.pptx")
        
        if os.path.exists(output_path):
            logger.info(f"Arquivo de saída encontrado: {output_path}")
            return jsonify({
                'status': 'completed',
                'progress': 100,
                'message': 'Apresentação criada com sucesso!',
                'redirectUrl': url_for('result', filename=secure_name)
            })
        
        # Se o processamento ainda está em andamento, calcular progresso
        file_ctime = os.path.getctime(file_path)
        time_elapsed = time.time() - file_ctime
        
        # Simula progresso baseado no tempo desde a criação do arquivo (25 segundos para completar)
        # Usamos 25 segundos porque o processamento real leva 20 segundos
        max_process_time = 25  # segundos
        progress = min(95, int((time_elapsed / max_process_time) * 100))
        
        # Determinar em qual etapa do processo estamos
        status = 'processing'
        if progress < 20:
            message = 'Extraindo texto do documento...'
        elif progress < 40:
            message = 'Analisando estrutura do conteúdo...'
        elif progress < 60:
            message = 'Gerando conteúdo com IA...'
        elif progress < 75:
            message = 'Criando slides e formatação...'
        else:
            message = 'Finalizando apresentação...'
        
        logger.info(f"Status check for {filename}: {status} - {progress}% - {message}")
        
        return jsonify({
            'status': status,
            'progress': progress,
            'message': message
        })
    
    # Função para simular a geração de um arquivo PowerPoint
    def simulate_file_generation(filename, delay=10):
        """Simula a geração de um arquivo PowerPoint após o upload"""
        logger.info(f"Iniciando simulação de geração de arquivo PowerPoint: {filename}")
        
        # Caminhos dos arquivos
        upload_folder = app.config['UPLOAD_FOLDER'].replace('/', os.sep)
        upload_dir = os.path.join(app.root_path, upload_folder)
        input_file_path = os.path.join(upload_dir, secure_filename(filename))
        output_file_path = os.path.join(upload_dir, f"processed_{secure_filename(filename)}.pptx")
        
        # Simula um tempo de processamento
        time.sleep(delay)
        
        # Para simulação, vamos apenas copiar o arquivo de entrada para o arquivo de saída
        try:
            shutil.copyfile(input_file_path, output_file_path)
            logger.info(f"Arquivo PowerPoint gerado com sucesso: {output_file_path}")
        except Exception as e:
            logger.error(f"Erro ao gerar arquivo PowerPoint: {e}")
    
    # Exemplo de uso da função de simulação (remover ou comentar em produção)
    @app.route('/simulate/<filename>')
    def simulate(filename):
        """Rota para simular a geração de arquivo PowerPoint"""
        thread = threading.Thread(target=simulate_file_generation, args=(filename,))
        thread.start()
        return jsonify({'status': 'File generation started', 'filename': filename}), 202
        
    # Log directories
    upload_dir = os.path.join(app.root_path, app.config['UPLOAD_FOLDER'])
    output_dir = os.path.join(app.root_path, 'static', 'outputs')
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)
    logger.info(f"Upload directory configured: {upload_dir}")
    logger.info(f"Output directory configured: {output_dir}")
    
    return app

# Create the application
app = create_app()

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5000))
    host = os.getenv('HOST', '0.0.0.0')
    
    logger.info(f"Starting DocToPPT application on {host}:{port}")
    logger.info(f"Configuration: {app.config.get('FLASK_ENV', 'default')}")
    logger.info(f"Debug mode: {app.config.get('DEBUG', False)}")
    
    # Verificar e instalar pacotes necessários
    check_and_install_missing_packages()

    try:
        # Importar PyMuPDF para garantir que está disponível
        import fitz
        logger.info("PyMuPDF (fitz) está instalado e disponível")
    except ImportError:
        logger.warning("PyMuPDF (fitz) não está instalado. Tentando instalar...")
        try:
            import subprocess
            import sys
            subprocess.run([sys.executable, "-m", "pip", "install", "PyMuPDF>=1.22.0"], 
                          stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, check=True)
            logger.info("PyMuPDF instalado com sucesso!")
        except Exception as e:
            logger.error(f"Erro ao instalar PyMuPDF: {e}")
    
    app.run(
        host=host,
        port=port,
        debug=app.config.get('DEBUG', False)
    )
